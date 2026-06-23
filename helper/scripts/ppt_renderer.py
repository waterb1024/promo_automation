"""
PPTX → 슬라이드별 PNG 변환.

흐름:
  1. PPTX 바이트를 임시 파일로 저장
  2. LibreOffice (soffice) 로 PPTX → PDF 변환
  3. poppler (pdf2image) 로 PDF 페이지별 PNG 추출
"""
from __future__ import annotations

import io
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import List, Dict


# LibreOffice 실행 파일을 PATH 또는 표준 위치에서 탐색
SOFFICE_CANDIDATES = [
    "soffice",
    "/opt/homebrew/bin/soffice",       # Apple Silicon Homebrew
    "/usr/local/bin/soffice",          # Intel Homebrew
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
]

# Poppler (pdftoppm) 가 설치된 디렉토리 후보
# launchd 환경의 PATH 가 brew 경로를 누락하는 경우 명시적으로 전달
POPPLER_PATH_CANDIDATES = [
    "/opt/homebrew/bin",   # Apple Silicon Homebrew
    "/usr/local/bin",      # Intel Homebrew
]


def _find_poppler_path():
    """pdftoppm 이 있는 디렉토리 반환. 못 찾으면 None (PATH 의존)."""
    for p in POPPLER_PATH_CANDIDATES:
        if Path(p).joinpath("pdftoppm").exists():
            return p
    return None


class PPTRenderError(RuntimeError):
    pass


def _find_soffice() -> str:
    """soffice 실행 파일 경로 반환. 없으면 PPTRenderError."""
    for cand in SOFFICE_CANDIDATES:
        # which 로 찾거나 직접 파일 존재 확인
        if "/" in cand:
            if Path(cand).exists():
                return cand
        else:
            found = shutil.which(cand)
            if found:
                return found
    raise PPTRenderError(
        "LibreOffice (soffice) 를 찾을 수 없습니다. "
        "설치: brew install --cask libreoffice"
    )


def _iter_text_shapes(shapes):
    """
    slide.shapes 를 재귀적으로 순회해 텍스트 프레임이 있는 leaf shape 만 yield.
    GROUP shape 은 자식으로 내려가고, 자기 자신은 yield 하지 않음.
    """
    for shape in shapes:
        # MSO_SHAPE_TYPE.GROUP == 6
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            yield from _iter_text_shapes(shape.shapes)
        elif getattr(shape, "has_text_frame", False):
            yield shape


def pptx_extract_texts(pptx_bytes: bytes) -> List[List[str]]:
    """
    PPTX 슬라이드별 텍스트 리스트 추출.

    추출 순서: 각 네모(shape)를 개별 단위로, 슬라이드 중앙 기준으로
              좌측 열 위→아래, 우측 열 위→아래 순서로 읽음.
    GROUP 안 shape 도 재귀적으로 포함.

    반환: [["슬라이드1 텍스트1", "슬라이드1 텍스트2"], ["슬라이드2 텍스트1", ...], ...]
    """
    try:
        from pptx import Presentation
    except ImportError:
        return []

    p = Presentation(io.BytesIO(pptx_bytes))
    slide_w = p.slide_width or 0
    mid_x = slide_w / 2 if slide_w else 0

    result = []
    for slide in p.slides:
        items = []  # (column, top, left, shape)
        for shape in _iter_text_shapes(slide.shapes):
            left = getattr(shape, "left", None) or 0
            top  = getattr(shape, "top",  None) or 0
            width = getattr(shape, "width", None) or 0
            cx = left + width / 2 if width else left
            column = 0 if cx < mid_x else 1   # 0 = 좌측, 1 = 우측
            items.append((column, top, left, shape))

        # 좌측 열 위→아래, 우측 열 위→아래
        items.sort(key=lambda x: (x[0], x[1], x[2]))

        texts = []
        for _, _, _, shape in items:
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    texts.append(t)
        result.append(texts)
    return result


def pptx_extract_texts_grouped(pptx_bytes: bytes) -> List[List[Dict]]:
    """
    PPTX 슬라이드별, shape(네모) 단위로 텍스트를 추출해 반환.

    각 텍스트 shape 을 독립적인 영역으로 취급 (병합 없음).
    GROUP 안 shape 도 재귀적으로 포함.

    정렬: 슬라이드 중앙 기준 좌측 열 위→아래, 우측 열 위→아래.
    영역 내 텍스트(paragraph 여러 줄): 위→아래 순서 유지.

    반환:
      [
        [  # 슬라이드 1
          {
            "bbox": [left, top, right, bottom],   # EMU
            "texts": ["...", "..."],              # shape 내 텍스트 (단락 순)
          },
          {...},
        ],
        [...],  # 슬라이드 2
      ]
    """
    try:
        from pptx import Presentation
    except ImportError:
        return []

    p = Presentation(io.BytesIO(pptx_bytes))
    slide_w = p.slide_width or 12192000

    result = []
    for slide in p.slides:
        mid_x = slide_w / 2
        items = []  # (column, top, left, right, bottom, texts)

        for shape in _iter_text_shapes(slide.shapes):
            left   = getattr(shape, "left",   None) or 0
            top    = getattr(shape, "top",    None) or 0
            width  = getattr(shape, "width",  None) or 0
            height = getattr(shape, "height", None) or 0
            right  = left + width
            bottom = top + height
            cx = left + width / 2 if width else left
            column = 0 if cx < mid_x else 1

            texts = []
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    texts.append(t)

            if not texts:
                continue

            items.append((column, top, left, right, bottom, texts))

        # 좌측 열 위→아래, 우측 열 위→아래
        items.sort(key=lambda x: (x[0], x[1], x[2]))

        group_list = [
            {
                "bbox": [left, top, right, bottom],
                "texts": texts,
            }
            for _, top, left, right, bottom, texts in items
        ]

        result.append(group_list)

    return result


def pptx_to_slide_pngs(pptx_bytes: bytes, dpi: int = 150) -> List[Dict]:
    """
    PPTX 바이트를 받아 슬라이드별 PNG bytes 리스트로 반환.

    반환:
      [{"index": 1, "bytes": b"...", "width": int, "height": int}, ...]
    """
    soffice = _find_soffice()

    try:
        from pdf2image import convert_from_path
    except ImportError as e:
        raise PPTRenderError(
            "pdf2image 가 설치되지 않았습니다. requirements.txt 확인 + install.sh 재실행."
        ) from e

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir)
        pptx_path = tmp / "input.pptx"
        pptx_path.write_bytes(pptx_bytes)

        # 1) PPTX → PDF
        try:
            proc = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf",
                 str(pptx_path), "--outdir", str(tmp)],
                capture_output=True,
                text=True,
                timeout=90,
            )
        except subprocess.TimeoutExpired as e:
            raise PPTRenderError(f"LibreOffice 변환 시간 초과: {e}") from e
        if proc.returncode != 0:
            raise PPTRenderError(
                f"LibreOffice 변환 실패 (rc={proc.returncode}): "
                f"{proc.stderr[:300] or proc.stdout[:300]}"
            )

        pdf_path = tmp / "input.pdf"
        if not pdf_path.exists():
            raise PPTRenderError("LibreOffice 변환은 성공했으나 PDF 가 생성되지 않음.")

        # 2) PDF → 페이지별 PIL Image (poppler 경로 launchd 안전하게 명시)
        poppler_path = _find_poppler_path()
        try:
            kwargs = {"dpi": dpi}
            if poppler_path:
                kwargs["poppler_path"] = poppler_path
            images = convert_from_path(str(pdf_path), **kwargs)
        except Exception as e:
            raise PPTRenderError(
                f"PDF → PNG 변환 실패 (poppler 누락 가능): {e}\n"
                f"설치 확인: brew install poppler"
            ) from e

        # 3) PIL → PNG bytes
        result = []
        for i, img in enumerate(images, 1):
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            result.append({
                "index": i,
                "bytes": buf.getvalue(),
                "width": img.width,
                "height": img.height,
            })
        return result
